#!/usr/bin/env python3
"""
PowerPoint theo đúng layout / mục lục trong Huong-Dan-Slide-KHONG-CHIA-SE.pdf:
  - Slide tiêu đề + thông tin sinh viên / GVHD
  - Danh sách nội dung (mục 1–14)
  - Mỗi mục: 1 slide "Title & Content" — toàn bộ gạch đầu dòng copy từ PDF
  - Tiếp theo: chỉ slide ảnh (Title Only + hình docs/diagram)

  pip install python-pptx
  python scripts/generate_thesis_pptx.py

  File ra: docs/Thesis-FoxGo-BaoCao.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Cần cài: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

REPO = Path(__file__).resolve().parents[1]
DIAGRAM = REPO / "docs" / "diagram"
OUT = REPO / "docs" / "Thesis-FoxGo-BaoCao.pptx"

# Kích thước slide (inch) — 10×7.5 mặc định
SLIDE_W_IN = 10.0
SLIDE_H_IN = 7.5
MARGIN_X = 0.45
TITLE_TOP_IN = 0.18
TITLE_H_IN = 0.58
GAP_IN = 0.12
MAX_IMG_W = SLIDE_W_IN - 2 * MARGIN_X
MAX_IMG_H = SLIDE_H_IN - TITLE_TOP_IN - TITLE_H_IN - GAP_IN - MARGIN_X
IMG_TOP_IN = TITLE_TOP_IN + TITLE_H_IN + GAP_IN


def _image_pixel_size(path: Path) -> tuple[int, int] | None:
    """Đọc kích thước pixel (PNG header; JPEG/ khác qua Pillow nếu có)."""
    try:
        head = path.read_bytes()[:32]
        if len(head) >= 24 and head[:8] == b"\x89PNG\r\n\x1a\n":
            w = int.from_bytes(head[16:20], "big")
            h = int.from_bytes(head[20:24], "big")
            if w > 0 and h > 0:
                return w, h
    except OSError:
        return None
    try:
        from PIL import Image  # type: ignore

        with Image.open(path) as im:
            return im.size
    except Exception:
        return None


def _fit_size_inches(path: Path) -> tuple[float, float]:
    """(width_in, height_in) vừa khít MAX_IMG_W × MAX_IMG_H, giữ tỷ lệ."""
    dim = _image_pixel_size(path)
    if not dim:
        return MAX_IMG_W * 0.95, MAX_IMG_H * 0.95
    pw, ph = dim
    if pw <= 0 or ph <= 0:
        return MAX_IMG_W * 0.95, MAX_IMG_H * 0.95
    ar = pw / ph
    box_ar = MAX_IMG_W / MAX_IMG_H
    if ar > box_ar:
        w_in = MAX_IMG_W
        h_in = MAX_IMG_W / ar
    else:
        h_in = MAX_IMG_H
        w_in = MAX_IMG_H * ar
    return w_in, h_in


# --- Đúng wording mục lục + gạch đầu dòng từ PDF hướng dẫn slide ---
TOC_LINES = [
    "1. Đặt vấn đề (Problem and scope)",
    "2. Thu thập yêu cầu và khảo sát nghiệp vụ (Requirement and survey)",
    "3. Phân tích và thiết kế nghiệp vụ (Domain or business model / Domain driven design)",
    "4. Phân tích và thiết kế chức năng của hệ thống (Use case diagram)",
    "5. Thiết kế kiến trúc hệ thống (System architecture design - SAD)",
    "6. (tùy chọn) Mẫu thiết kế phần mềm (Design Pattern)",
    "7. (tùy chọn) Cài đặt và triển khai (CI/CD)",
    "8. (tùy chọn) Ứng dụng AI (AI / AI agent pipeline)",
    "9. (tùy chọn) Giám sát hệ thống (System monitor)",
    "10. (tùy chọn) Tối ưu kiến trúc giao dịch (Transaction process) — thanh toán trực tuyến MSA",
    "11. (tùy chọn) Tối ưu hóa hệ thống đáp ứng nhu cầu (System scalability)",
    "12. (tùy chọn) Tăng cường bảo mật hệ thống và dữ liệu (Information Security, Data security)",
    "13. Kết luận",
    "14. Demo",
]

SECTIONS: list[dict] = [
    {
        "slide_title": "1. Đặt vấn đề (Problem and scope)",
        "subtitle": "(fresher level)",
        "bullets": [
            "Mô tả bài toán và phạm vi ngắn gọn không nên đọc slide, nên dùng "
            "danh sách gạch đầu dòng liệt kê ý chính và tự thuyết minh chi tiết.",
        ],
        "images": ["overview_arch_system.png"],
    },
    {
        "slide_title": "2. Thu thập yêu cầu và khảo sát nghiệp vụ (Requirement and survey)",
        "subtitle": "(fresher level)",
        "bullets": [
            "Thu thập yêu cầu từ khách hàng.",
            "Khảo sát các hệ thống thực tế điển hình: nên dùng hình ảnh biểu đồ phân "
            "tích dữ liệu và tự thuyết minh diễn giải — từ đó thể hiện kinh nghiệm người làm dự án.",
            "Cách đọc biểu đồ phân tích.",
            "Câu hỏi muốn giải quyết là gì?",
            "Số liệu biểu đồ cho thấy thông tin gì?",
            "Rút ra kết luận gì áp dụng cho dự án?",
        ],
        "images": ["uc_journey_roles_en.png"],
    },
    {
        "slide_title": "3. Phân tích và thiết kế nghiệp vụ (DDD)",
        "subtitle": "(fresher level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
        ],
        "images": [
            "ddd_subdomain_map.png",
            "ddd_bc_service_map.png",
            "erd_core_bounded_contexts.png",
        ],
    },
    {
        "slide_title": "4. Phân tích và thiết kế chức năng (Use case diagram)",
        "subtitle": "(junior level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
        ],
        "images": ["bpmn_booking_swimlanes.png", "uc_journey_roles_en.png"],
    },
    {
        "slide_title": "5. Thiết kế kiến trúc hệ thống (SAD)",
        "subtitle": "(middle level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
        ],
        "images": [
            "arch_layers_services_infra.png",
            "arch_ms_port_map.png",
            "arch_ms_comm_flow.png",
        ],
    },
    {
        "slide_title": "6. (tùy chọn) Mẫu thiết kế phần mềm (Design Pattern)",
        "subtitle": "(middle level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
        ],
        "images": ["dp_pattern_catalog.png", "arch_integration_pattern_map.png"],
    },
    {
        "slide_title": "7. (tùy chọn) Cài đặt và triển khai (CI/CD)",
        "subtitle": "(middle level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
        ],
        "images": [
            "cicd_github_actions_docker.png",
            "deploy_docker_swarm_topology.png",
            "deploy_swarm_aws_asbuilt.png",
        ],
    },
    {
        "slide_title": "8. (tùy chọn) Ứng dụng AI (AI / AI agent pipeline)",
        "subtitle": "(middle level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
        ],
        "images": [
            "ai_agent_rag_retrieval.png",
            "ai_ml_sklearn_train_infer_pipeline.png",
        ],
    },
    {
        "slide_title": "9. (tùy chọn) Giám sát hệ thống (System monitor)",
        "subtitle": "(middle level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
            "ELK (Elastic Logstash Kibana): Logging.",
            "Grafana / Prometheus: Performance.",
        ],
        "images": ["mon_prometheus_grafana_stack.png"],
    },
    {
        "slide_title": "10. (tùy chọn) Tối ưu kiến trúc giao dịch (Transaction process — MSA)",
        "subtitle": "(senior level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
        ],
        "images": [
            "seq_payment_after_ride_complete.png",
            "act_payment_online_saga.png",
            "seq_payment_ipn_idempotent.png",
        ],
    },
    {
        "slide_title": "11. (tùy chọn) Tối ưu hóa hệ thống đáp ứng nhu cầu (System scalability)",
        "subtitle": "(senior level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
            "Đánh giá hiệu năng mở hệ thống.",
            "Ví dụ: quy mô lớn — Yêu cầu hệ thống đáp ứng nhanh (SLA).",
        ],
        "images": ["arch_scalability_approach.png"],
    },
    {
        "slide_title": "12. (tùy chọn) Tăng cường bảo mật hệ thống và dữ liệu",
        "subtitle": "(senior level)",
        "bullets": [
            "Nên dùng hình ảnh tự vẽ theo hệ thống của dự án và tự thuyết minh diễn giải.",
        ],
        "images": ["sec_trust_boundary_multitier.png"],
    },
    {
        "slide_title": "13. Kết luận",
        "subtitle": "",
        "bullets": [
            "Mô tả bài toán ngắn gọn — danh sách gạch đầu dòng ý chính; tự thuyết minh chi tiết.",
            "Dùng bảng so sánh ưu điểm / nhược điểm.",
        ],
        "images": [],
    },
    {
        "slide_title": "14. Demo",
        "subtitle": "",
        "bullets": [
            "Bắt đầu demo ứng dụng trực tiếp.",
            "Nạp dữ liệu mẫu trước; chọn quy trình hoàn chỉnh; dành thời gian cho chức năng đặc sắc (đặc biệt AI).",
        ],
        "images": ["act_booking_end_to_end.png"],
    },
]


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    ttf = slide.shapes.title.text_frame
    ttf.text = title
    ttf.word_wrap = True
    for p in ttf.paragraphs:
        p.font.size = Pt(28) if len(title) < 40 else Pt(22)
    ph = slide.placeholders[1]
    if ph:
        ph.text = subtitle
        stf = ph.text_frame
        stf.word_wrap = True
        for p in stf.paragraphs:
            p.font.size = Pt(14)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    ttf = slide.shapes.title.text_frame
    ttf.text = title
    ttf.word_wrap = True
    for p in ttf.paragraphs:
        p.font.size = Pt(20) if len(title) < 72 else Pt(14)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.word_wrap = True
    body.margin_bottom = Inches(0.08)
    body.margin_left = Inches(0.05)
    for i, line in enumerate(bullets):
        if i == 0:
            body.text = line
        else:
            p = body.add_paragraph()
            p.text = line
            p.level = 0
    for p in body.paragraphs:
        p.font.size = Pt(11)
        p.space_after = Pt(3)


def add_image_slide(prs: Presentation, title: str, image_name: str) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    path = DIAGRAM / image_name

    title_w = Inches(SLIDE_W_IN - 2 * MARGIN_X)
    title_box = slide.shapes.add_textbox(
        Inches(MARGIN_X), Inches(TITLE_TOP_IN), title_w, Inches(TITLE_H_IN)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title
    for p in tf.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True

    if not path.is_file():
        miss = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(3), title_w, Inches(0.8))
        miss.text_frame.text = f"(Thiếu ảnh: docs/diagram/{image_name})"
        return

    w_in, h_in = _fit_size_inches(path)
    left = (SLIDE_W_IN - w_in) / 2
    zone_bottom = SLIDE_H_IN - MARGIN_X
    zone_h = zone_bottom - IMG_TOP_IN
    top = IMG_TOP_IN + max(0.0, (zone_h - h_in) / 2)

    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(w_in), height=Inches(h_in))


def main() -> None:
    prs = Presentation()
    prs.slide_width = int(Inches(10))
    prs.slide_height = int(Inches(7.5))

    # --- Slide bìa theo PDF ---
    add_title_slide(
        prs,
        "TIÊU ĐỀ",
        "Thông tin sinh viên:\n"
        "Người nói đầu tiên: …\n"
        "Người nói cuối cùng: …\n"
        "GVHD: ThS Huỳnh Nam",
    )

    add_bullet_slide(prs, "Danh sách nội dung", TOC_LINES)

    for sec in SECTIONS:
        title = sec["slide_title"]
        sub = sec.get("subtitle") or ""
        bullets = list(sec["bullets"])
        if sub:
            bullets = [sub, *bullets]
        add_bullet_slide(prs, title, bullets)
        for img in sec.get("images") or []:
            short = Path(img).stem.replace("_", " ")
            add_image_slide(prs, f"Minh họa: {short}", img)

    add_title_slide(prs, "Cảm ơn thầy cô và hội đồng", "Hỏi & đáp")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Đã tạo: {OUT}")


if __name__ == "__main__":
    main()
